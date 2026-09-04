"""
RAG Document Processor — extract, chunk, embed, store.
"""
import logging
import uuid
from typing import Optional
from sqlalchemy.orm import Session

from app.core.config import settings

logger = logging.getLogger(__name__)


def _extract_text(file_path: str, file_ext: str) -> str:
    """Extract plain text from a document."""
    try:
        if file_ext == "pdf":
            from pypdf import PdfReader
            reader = PdfReader(file_path)
            return "\n".join(page.extract_text() or "" for page in reader.pages)
        elif file_ext == "docx":
            from docx import Document
            doc = Document(file_path)
            return "\n".join(p.text for p in doc.paragraphs)
        elif file_ext == "txt":
            with open(file_path, "r", encoding="utf-8", errors="replace") as f:
                return f.read()
        elif file_ext == "pptx":
            from pptx import Presentation
            prs = Presentation(file_path)
            texts = []
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        texts.append(shape.text)
            return "\n".join(texts)
    except Exception as e:
        logger.error(f"Text extraction error for {file_path}: {e}")
    return ""


def _chunk_text(text: str, chunk_size: int = 800, overlap: int = 100) -> list[str]:
    """Split text into overlapping chunks."""
    if not text.strip():
        return []
    chunks = []
    start = 0
    while start < len(text):
        end = start + chunk_size
        chunk = text[start:end].strip()
        if chunk:
            chunks.append(chunk)
        start = end - overlap
        if start >= len(text):
            break
    return chunks


class DocumentProcessor:
    def __init__(self):
        self._chroma_client = None
        self._embedding_fn = None

    def _get_chroma(self):
        if self._chroma_client is None:
            import chromadb
            self._chroma_client = chromadb.PersistentClient(
                path=settings.CHROMA_PERSIST_DIR
            )
        return self._chroma_client

    def _get_collection(self, user_id: str):
        """Get or create a ChromaDB collection per user."""
        client = self._get_chroma()
        collection_name = f"user_{user_id.replace('-', '_')}"
        return client.get_or_create_collection(
            name=collection_name,
            metadata={"hnsw:space": "cosine"},
        )

    async def process(
        self,
        doc_id: str,
        file_path: str,
        file_ext: str,
        user_id: str,
        db: Session,
    ) -> bool:
        """Full RAG pipeline: extract → chunk → embed → store."""
        from app.models.documents import Document, DocumentChunk

        text = _extract_text(file_path, file_ext)
        if not text.strip():
            logger.warning(f"No text extracted from {file_path}")
            return False

        chunks = _chunk_text(text)
        if not chunks:
            return False

        collection = self._get_collection(user_id)
        doc = db.query(Document).filter(Document.id == doc_id).first()
        if not doc:
            return False

        chunk_ids = []
        chunk_texts = []
        chunk_metas = []

        db_chunks = []
        for i, chunk in enumerate(chunks):
            chroma_id = f"{doc_id}_chunk_{i}"
            chunk_ids.append(chroma_id)
            chunk_texts.append(chunk)
            chunk_metas.append({
                "doc_id": str(doc_id),
                "doc_name": doc.original_filename,
                "chunk_index": i,
                "subject": doc.subject or "",
                "doc_type": doc.document_type or "",
            })
            db_chunks.append(DocumentChunk(
                document_id=doc_id,
                chunk_index=i,
                content=chunk,
                chroma_id=chroma_id,
            ))

        # Add to ChromaDB (uses built-in embedding model)
        collection.add(
            ids=chunk_ids,
            documents=chunk_texts,
            metadatas=chunk_metas,
        )

        # Save chunks to DB
        for chunk_obj in db_chunks:
            db.add(chunk_obj)

        doc.is_processed = True
        doc.chunk_count = len(chunks)
        doc.chroma_collection_id = f"user_{user_id.replace('-', '_')}"
        db.commit()

        logger.info(f"Processed document {doc_id}: {len(chunks)} chunks")
        return True

    async def query(
        self,
        user_id: str,
        query: str,
        document_ids: Optional[list] = None,
        n_results: int = 5,
    ) -> list[dict]:
        """Query ChromaDB for relevant chunks."""
        try:
            collection = self._get_collection(user_id)
            where = None
            if document_ids:
                where = {"doc_id": {"$in": [str(d) for d in document_ids]}}

            results = collection.query(
                query_texts=[query],
                n_results=n_results,
                where=where,
            )
            if not results or not results["documents"]:
                return []

            chunks = []
            for i, doc_text in enumerate(results["documents"][0]):
                meta = results["metadatas"][0][i] if results["metadatas"] else {}
                chunks.append({
                    "content": doc_text,
                    "doc_name": meta.get("doc_name", "Unknown"),
                    "doc_id": meta.get("doc_id", ""),
                    "chunk_index": meta.get("chunk_index", 0),
                })
            return chunks
        except Exception as e:
            logger.error(f"RAG query error: {e}")
            return []

    async def delete_document(self, doc_id: str, user_id: str):
        """Remove document chunks from ChromaDB."""
        try:
            collection = self._get_collection(user_id)
            # Get all chunk IDs for this document
            results = collection.get(where={"doc_id": str(doc_id)})
            if results and results["ids"]:
                collection.delete(ids=results["ids"])
        except Exception as e:
            logger.error(f"ChromaDB delete error: {e}")


# Singleton
document_processor = DocumentProcessor()
